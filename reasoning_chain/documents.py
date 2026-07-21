"""Persistent multi-format extraction and bounded document retrieval for session RAG."""

from __future__ import annotations

import csv
import io
import math
import os
import re
import uuid
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime, timezone
from html import escape
from pathlib import Path
from typing import Optional

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel, ConfigDict

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover - guarded by the runtime dependency
    PdfReader = None


def _env_int(name: str, default: int, minimum: int = 1) -> int:
    try:
        return max(minimum, int(os.environ.get(name, default)))
    except (TypeError, ValueError):
        return default


@dataclass(frozen=True)
class DocumentSettings:
    max_file_bytes: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_MAX_FILE_BYTES",
            _env_int("PDF_MAX_FILE_BYTES", 10 * 1024 * 1024),
        )
    )
    max_pages: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_MAX_UNITS", _env_int("PDF_MAX_PAGES", 200)
        )
    )
    max_extracted_chars: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_MAX_EXTRACTED_CHARS",
            _env_int("PDF_MAX_EXTRACTED_CHARS", 2_000_000),
        )
    )
    chunk_words: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_CHUNK_WORDS", _env_int("PDF_CHUNK_WORDS", 450, 100), 100
        )
    )
    chunk_overlap_words: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_CHUNK_OVERLAP_WORDS",
            _env_int("PDF_CHUNK_OVERLAP_WORDS", 75),
        )
    )
    retrieval_chunks: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_RETRIEVAL_CHUNKS", _env_int("PDF_RETRIEVAL_CHUNKS", 8)
        )
    )
    retrieval_max_chars: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_RETRIEVAL_MAX_CHARS",
            _env_int("PDF_RETRIEVAL_MAX_CHARS", 14_000),
        )
    )
    max_documents_per_session: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_MAX_PER_SESSION",
            _env_int("PDF_MAX_DOCUMENTS_PER_SESSION", 10),
        )
    )
    ttl_seconds: int = field(
        default_factory=lambda: _env_int(
            "DOCUMENT_TTL_SECONDS", _env_int("PDF_TTL_SECONDS", 2_592_000)
        )
    )


class DocumentError(ValueError):
    """A safe error that can be returned to a document API caller."""


class DocumentMetadata(BaseModel):
    model_config = ConfigDict(extra="forbid")

    document_id: str
    session_id: str
    filename: str
    document_type: str = "pdf"
    unit_label: str = "page"
    unit_count: int = 0
    status: str = "ready"
    page_count: int
    chunk_count: int
    extracted_chars: int
    created_at: str


class DocumentChunk(BaseModel):
    model_config = ConfigDict(extra="forbid")

    chunk_id: str
    document_id: str
    filename: str
    page: Optional[int] = None
    location: str = ""
    text: str

    @property
    def citation(self) -> str:
        locator = self.location or f"page {self.page or 1}"
        return f"[{self.filename}, {locator}]"


@dataclass(frozen=True)
class DocumentRetrieval:
    context: str = ""
    citations: list[str] = field(default_factory=list)
    document_ids: list[str] = field(default_factory=list)
    chunk_count: int = 0


GENERIC_DOCUMENT_TERMS = {
    "about",
    "document",
    "explain",
    "file",
    "isme",
    "iska",
    "iske",
    "iski",
    "isko",
    "batao",
    "pdf",
    "summarise",
    "summarize",
    "summary",
    "samjhao",
    "this",
    "usme",
}
SUPPORTED_DOCUMENT_TYPES = {
    ".csv": ("csv", "row"),
    ".docx": ("docx", "block"),
    ".md": ("markdown", "section"),
    ".pdf": ("pdf", "page"),
    ".pptx": ("pptx", "slide"),
    ".tsv": ("tsv", "row"),
    ".txt": ("text", "section"),
    ".xlsx": ("xlsx", "sheet"),
}
STOP_TERMS = {
    "a",
    "an",
    "and",
    "are",
    "for",
    "from",
    "h",
    "hai",
    "in",
    "is",
    "ka",
    "ke",
    "ki",
    "ko",
    "kya",
    "me",
    "mein",
    "of",
    "on",
    "the",
    "to",
    "what",
}


def _safe_filename(filename: Optional[str]) -> str:
    name = Path(filename or "document.txt").name.strip()
    name = re.sub(r"[\x00-\x1f\x7f]", "", name)[:120]
    extension = Path(name).suffix.lower()
    if extension not in SUPPORTED_DOCUMENT_TYPES:
        supported = ", ".join(sorted(SUPPORTED_DOCUMENT_TYPES))
        raise DocumentError(f"unsupported document type; allowed extensions: {supported}")
    return name or "document.txt"


def _normalize_page_text(value: str) -> str:
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in value.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _chunk_unit(
    text: str,
    location: str,
    document_id: str,
    filename: str,
    settings: DocumentSettings,
    page: Optional[int] = None,
) -> list[DocumentChunk]:
    words = text.split()
    if not words:
        return []
    chunk_words = settings.chunk_words
    overlap = min(settings.chunk_overlap_words, chunk_words // 3)
    step = max(1, chunk_words - overlap)
    chunks = []
    for index, start in enumerate(range(0, len(words), step), start=1):
        segment = words[start : start + chunk_words]
        if not segment:
            break
        chunks.append(
            DocumentChunk(
                chunk_id=f"{document_id}:p{page}:c{index}",
                document_id=document_id,
                filename=filename,
                page=page,
                location=location,
                text=" ".join(segment),
            )
        )
        if start + chunk_words >= len(words):
            break
    return chunks


def _extract_pdf_units(
    data: bytes,
) -> tuple[list[tuple[str, str, Optional[int]]], int]:
    if not data.startswith(b"%PDF-"):
        raise DocumentError("file content is not a valid PDF")
    if PdfReader is None:
        raise DocumentError("PDF extraction dependency is not installed")
    try:
        reader = PdfReader(io.BytesIO(data), strict=False)
        if reader.is_encrypted and reader.decrypt("") == 0:
            raise DocumentError("password-protected PDFs are not supported")
        page_count = len(reader.pages)
    except DocumentError:
        raise
    except Exception as exc:
        raise DocumentError(f"unable to parse PDF: {exc}") from exc
    if page_count < 1:
        raise DocumentError("PDF has no pages")
    units = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            text = _normalize_page_text(page.extract_text() or "")
        except Exception:
            text = ""
        if text:
            units.append((f"page {page_number}", text, page_number))
    return units, page_count


def _decode_text(data: bytes) -> str:
    try:
        return data.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DocumentError("text documents must use UTF-8 encoding") from exc


def _extract_text_units(data: bytes) -> tuple[list[tuple[str, str, None]], int]:
    text = _normalize_page_text(_decode_text(data))
    return [("section 1", text, None)] if text else [], 1


def _extract_delimited_units(
    data: bytes, delimiter: str
) -> tuple[list[tuple[str, str, None]], int]:
    try:
        rows = list(csv.reader(io.StringIO(_decode_text(data)), delimiter=delimiter))
    except csv.Error as exc:
        raise DocumentError(f"unable to parse delimited document: {exc}") from exc
    units = []
    batch_size = 50
    for start in range(0, len(rows), batch_size):
        batch = rows[start : start + batch_size]
        rendered = "\n".join(" | ".join(str(cell).strip() for cell in row) for row in batch)
        end = start + len(batch)
        units.append((f"rows {start + 1}-{end}", rendered, None))
    return units, len(rows)


def _extract_docx_units(data: bytes) -> tuple[list[tuple[str, str, None]], int]:
    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise DocumentError(f"unable to parse DOCX: {exc}") from exc
    blocks = []
    for paragraph in document.paragraphs:
        text = _normalize_page_text(paragraph.text)
        if text:
            style = str(getattr(paragraph.style, "name", "") or "")
            blocks.append(f"{style}: {text}" if style.lower().startswith("heading") else text)
    for table_index, table in enumerate(document.tables, start=1):
        rows = [
            " | ".join(_normalize_page_text(cell.text) for cell in row.cells)
            for row in table.rows
        ]
        if rows:
            blocks.append(f"Table {table_index}\n" + "\n".join(rows))
    units = []
    batch_size = 20
    for start in range(0, len(blocks), batch_size):
        batch = blocks[start : start + batch_size]
        units.append(
            (
                f"blocks {start + 1}-{start + len(batch)}",
                "\n".join(batch),
                None,
            )
        )
    return units, len(blocks)


def _cell_text(value) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()
        except Exception:
            pass
    return str(value)


def _extract_xlsx_units(data: bytes) -> tuple[list[tuple[str, str, None]], int]:
    try:
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=False)
    except Exception as exc:
        raise DocumentError(f"unable to parse XLSX: {exc}") from exc
    units = []
    try:
        for worksheet in workbook.worksheets:
            batch = []
            batch_start = 1
            for row_number, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = [_cell_text(value).strip() for value in row]
                if not any(values):
                    continue
                if not batch:
                    batch_start = row_number
                batch.append(" | ".join(values))
                if len(batch) == 50:
                    units.append(
                        (
                            f'sheet "{worksheet.title}", rows {batch_start}-{row_number}',
                            "\n".join(batch),
                            None,
                        )
                    )
                    batch = []
            if batch:
                end = batch_start + len(batch) - 1
                units.append(
                    (
                        f'sheet "{worksheet.title}", rows {batch_start}-{end}',
                        "\n".join(batch),
                        None,
                    )
                )
        return units, len(workbook.sheetnames)
    finally:
        workbook.close()


def _extract_pptx_units(data: bytes) -> tuple[list[tuple[str, str, None]], int]:
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise DocumentError(f"unable to parse PPTX: {exc}") from exc
    units = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = _normalize_page_text(getattr(shape, "text", "") or "")
            if text:
                parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(_normalize_page_text(cell.text) for cell in row.cells))
        if parts:
            units.append((f"slide {slide_number}", "\n".join(parts), None))
    return units, len(presentation.slides)


def extract_document(
    data: bytes,
    filename: Optional[str],
    settings: Optional[DocumentSettings] = None,
) -> tuple[DocumentMetadata, list[DocumentChunk]]:
    """Validate and extract a supported document into locator-preserving chunks."""
    settings = settings or DocumentSettings()
    safe_name = _safe_filename(filename)
    if len(data) > settings.max_file_bytes:
        raise DocumentError(
            f"document exceeds the {settings.max_file_bytes // (1024 * 1024)} MB limit"
        )
    extension = Path(safe_name).suffix.lower()
    document_type, unit_label = SUPPORTED_DOCUMENT_TYPES[extension]
    if extension == ".pdf":
        units, unit_count = _extract_pdf_units(data)
    elif extension == ".docx":
        units, unit_count = _extract_docx_units(data)
    elif extension in {".txt", ".md"}:
        units, unit_count = _extract_text_units(data)
    elif extension in {".csv", ".tsv"}:
        units, unit_count = _extract_delimited_units(
            data, "," if extension == ".csv" else "\t"
        )
    elif extension == ".xlsx":
        units, unit_count = _extract_xlsx_units(data)
    else:
        units, unit_count = _extract_pptx_units(data)
    unit_limit = (
        settings.max_pages
        if unit_label in {"page", "sheet", "slide"}
        else settings.max_pages * 50
    )
    if unit_count > unit_limit:
        raise DocumentError(
            f"document exceeds the {unit_limit}-{unit_label} processing limit"
        )

    extracted_chars = sum(len(text) for _, text, _ in units)
    if extracted_chars > settings.max_extracted_chars:
        raise DocumentError("document extracted text exceeds the configured safety limit")
    document_id = str(uuid.uuid4())
    chunks = []
    for location, text, page in units:
        chunks.extend(
            _chunk_unit(
                text,
                location,
                document_id,
                safe_name,
                settings,
                page=page,
            )
        )

    if extracted_chars < 20 or not chunks:
        raise DocumentError(
            "document has no extractable text; scanned/image documents require OCR, which is not "
            "configured"
        )
    metadata = DocumentMetadata(
        document_id=document_id,
        session_id="",
        filename=safe_name,
        document_type=document_type,
        unit_label=unit_label,
        unit_count=unit_count,
        page_count=unit_count,
        chunk_count=len(chunks),
        extracted_chars=extracted_chars,
        created_at=datetime.now(timezone.utc).isoformat(),
    )
    return metadata, chunks


def extract_pdf(
    data: bytes,
    filename: Optional[str],
    settings: Optional[DocumentSettings] = None,
) -> tuple[DocumentMetadata, list[DocumentChunk]]:
    """Backward-compatible PDF extraction wrapper."""
    if Path(filename or "").suffix.lower() != ".pdf":
        raise DocumentError("only .pdf files are supported by extract_pdf")
    return extract_document(data, filename, settings)


def _tokens(value: str) -> list[str]:
    return [token for token in re.findall(r"[^\W_]+", value.lower(), re.UNICODE) if len(token) > 1]


class RedisDocumentStore:
    """Stores extracted chunks in Redis and retrieves only within one session."""

    def __init__(self, client, settings: Optional[DocumentSettings] = None):
        self.client = client
        self.settings = settings or DocumentSettings()

    @staticmethod
    def _session_key(session_id: str) -> str:
        return f"session:{session_id}:documents"

    @staticmethod
    def _metadata_key(document_id: str) -> str:
        return f"document:{document_id}:metadata"

    @staticmethod
    def _chunks_key(document_id: str) -> str:
        return f"document:{document_id}:chunks"

    def ingest(self, session_id: str, filename: Optional[str], data: bytes) -> DocumentMetadata:
        session_key = self._session_key(session_id)
        existing = self.client.lrange(session_key, 0, -1)
        if len(existing) >= self.settings.max_documents_per_session:
            raise DocumentError(
                f"session already has {self.settings.max_documents_per_session} documents"
            )
        metadata, chunks = extract_document(data, filename, self.settings)
        metadata.session_id = session_id
        metadata_key = self._metadata_key(metadata.document_id)
        chunks_key = self._chunks_key(metadata.document_id)
        pipe = self.client.pipeline()
        pipe.set(metadata_key, metadata.model_dump_json(), ex=self.settings.ttl_seconds)
        pipe.rpush(chunks_key, *(chunk.model_dump_json() for chunk in chunks))
        pipe.expire(chunks_key, self.settings.ttl_seconds)
        pipe.rpush(session_key, metadata.document_id)
        pipe.expire(session_key, self.settings.ttl_seconds)
        pipe.execute()
        return metadata

    def list(self, session_id: str) -> list[DocumentMetadata]:
        document_ids = self.client.lrange(self._session_key(session_id), 0, -1)
        documents = []
        for document_id in document_ids:
            raw = self.client.get(self._metadata_key(str(document_id)))
            if raw:
                try:
                    metadata = DocumentMetadata.model_validate_json(raw)
                except Exception:
                    continue
                if metadata.session_id == session_id:
                    documents.append(metadata)
        return documents

    def delete(self, session_id: str, document_id: str) -> bool:
        raw = self.client.get(self._metadata_key(document_id))
        if not raw:
            return False
        metadata = DocumentMetadata.model_validate_json(raw)
        if metadata.session_id != session_id:
            return False
        pipe = self.client.pipeline()
        pipe.delete(self._metadata_key(document_id), self._chunks_key(document_id))
        pipe.lrem(self._session_key(session_id), 0, document_id)
        pipe.execute()
        return True

    def _load_chunks(self, documents: list[DocumentMetadata]) -> list[DocumentChunk]:
        chunks = []
        for document in documents:
            raw_chunks = self.client.lrange(
                self._chunks_key(document.document_id), 0, -1
            )
            for raw in raw_chunks:
                try:
                    chunks.append(DocumentChunk.model_validate_json(raw))
                except Exception:
                    continue
        return chunks

    def retrieve(self, session_id: str, query: str) -> DocumentRetrieval:
        documents = self.list(session_id)
        chunks = self._load_chunks(documents)
        if not chunks:
            return DocumentRetrieval()

        raw_terms = set(_tokens(query))
        generic_intent = bool(raw_terms & GENERIC_DOCUMENT_TERMS)
        query_terms = raw_terms - STOP_TERMS - GENERIC_DOCUMENT_TERMS
        selected: list[DocumentChunk]
        if not query_terms:
            if not generic_intent:
                return DocumentRetrieval()
            selected = self._representative_chunks(chunks)
        else:
            selected = self._ranked_chunks(chunks, query_terms)
            if not selected and generic_intent:
                selected = self._representative_chunks(chunks)
        if not selected:
            return DocumentRetrieval()
        return self._format_retrieval(selected)

    def _representative_chunks(self, chunks: list[DocumentChunk]) -> list[DocumentChunk]:
        limit = min(self.settings.retrieval_chunks, len(chunks))
        if limit == len(chunks):
            return chunks
        indexes = {
            round(index * (len(chunks) - 1) / max(1, limit - 1)) for index in range(limit)
        }
        return [chunks[index] for index in sorted(indexes)]

    def _ranked_chunks(
        self, chunks: list[DocumentChunk], query_terms: set[str]
    ) -> list[DocumentChunk]:
        token_counts = [Counter(_tokens(chunk.text)) for chunk in chunks]
        document_frequency = {
            term: sum(1 for counts in token_counts if counts.get(term)) for term in query_terms
        }
        scored = []
        for index, (chunk, counts) in enumerate(zip(chunks, token_counts)):
            score = 0.0
            for term in query_terms:
                frequency = counts.get(term, 0)
                if frequency:
                    inverse_frequency = math.log(
                        (len(chunks) + 1) / (document_frequency[term] + 1)
                    ) + 1
                    score += (1 + math.log(frequency)) * inverse_frequency
                if term in chunk.filename.lower():
                    score += 1.5
            if score > 0:
                scored.append((score, -index, chunk))
        scored.sort(reverse=True, key=lambda item: (item[0], item[1]))
        return [item[2] for item in scored[: self.settings.retrieval_chunks]]

    def _format_retrieval(self, chunks: list[DocumentChunk]) -> DocumentRetrieval:
        included = []
        used_chars = 0
        for chunk in chunks:
            if included and used_chars + len(chunk.text) > self.settings.retrieval_max_chars:
                break
            included.append(chunk)
            used_chars += len(chunk.text)
        passages = []
        citations = []
        document_ids = []
        for chunk in included:
            citation = chunk.citation
            citations.append(citation)
            if chunk.document_id not in document_ids:
                document_ids.append(chunk.document_id)
            passages.append(
                f'<passage document_id="{escape(chunk.document_id, quote=True)}" '
                f'citation="{escape(citation, quote=True)}">'
                f"{escape(chunk.text, quote=False)}</passage>"
            )
        context = (
            '<document_context trust="untrusted" priority="critical">\n'
            + "\n".join(passages)
            + "\n</document_context>"
        )
        return DocumentRetrieval(
            context=context,
            citations=list(dict.fromkeys(citations)),
            document_ids=document_ids,
            chunk_count=len(included),
        )
